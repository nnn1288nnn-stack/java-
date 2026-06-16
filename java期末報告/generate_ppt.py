from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── Color palette ───
BG_DARK   = RGBColor(0x1E, 0x2A, 0x38)
BG_CARD   = RGBColor(0x27, 0x38, 0x4A)
ACCENT    = RGBColor(0x4E, 0xCB, 0x71)  # green
ACCENT2   = RGBColor(0x5B, 0x9B, 0xD5)  # blue
ACCENT3   = RGBColor(0xF5, 0xA6, 0x23)  # orange
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
SOFT_GREEN = RGBColor(0x81, 0xC7, 0x84)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from pptx.oxml.ns import qn
        spPr = shape._element.spPr
        spPr.set(qn('a:alpha'), str(int(alpha * 1000)))
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='微軟正黑體'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '微軟正黑體'
        p.space_after = spacing
    return txBox

def add_code_block(slide, left, top, width, height, code_lines, font_size=11, bg_color=RGBColor(0x15, 0x1E, 0x2A)):
    shape = add_shape(slide, left, top, width, height, bg_color)
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xE6, 0xDB, 0x74) if line.strip().startswith('//') else LIGHT_GRAY
        p.font.name = 'Consolas'
        p.space_after = Pt(1)
    return txBox

def section_header(slide, title, subtitle=""):
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.5), ACCENT)
    add_text_box(slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.9), title, font_size=32, bold=True, color=RGBColor(0x1E, 0x2A, 0x38))
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.5), subtitle, font_size=16, color=RGBColor(0x1E, 0x2A, 0x38))

def footer(slide, text="農場物語 — Java 2D 簡易農場遊戲"):
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(6), Inches(0.4), text, font_size=10, color=LIGHT_GRAY)

# ================================================================
# SLIDE 1 — Title
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(2.0), prs.slide_width, Inches(3.5), BG_CARD)
# Decorative top bar
add_shape(slide, Inches(0), Inches(1.7), prs.slide_width, Inches(0.08), ACCENT)
# Leaf icon placeholder
add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.2), "🌾 農場物語", font_size=48, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.8), "以 Java Swing/AWT 實作之 2D 簡易農場遊戲", font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.6), "靈感源自《星露谷物語》(Stardew Valley) 核心玩法", font_size=16, color=SOFT_GREEN, alignment=PP_ALIGN.CENTER)
footer(slide, "")

# ================================================================
# SLIDE 2 — 目錄
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "目錄 CONTENTS", "本報告導覽")

items = [
    "1. 專案概述",
    "2. 開發環境與工具",
    "3. 系統架構",
    "4. 核心類別說明",
    "5. 耕作流程與狀態機",
    "6. 操作說明",
    "7. 成果展示",
    "8. 未來展望",
]
add_bullet_slide(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(4.5), items, font_size=22, color=LIGHT_GRAY, spacing=Pt(12))
footer(slide)

# ================================================================
# SLIDE 3 — 專案概述
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "專案概述", "Project Overview")

add_text_box(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(5), "",
             font_size=14, color=LIGHT_GRAY)

items_left = [
    "▸ 使用 Java Swing / AWT 開發",
    "▸ 800 × 600 固定大小遊戲視窗",
    "▸ 60 FPS 固定幀率遊戲主迴圈",
    "▸ 25 × 19 網格地圖（32×32 px / 格）",
    "▸ WASD 網格移動 + 碰撞偵測",
    "▸ 完整農耕流程：翻土→播種→澆水→成長→收成",
    "▸ 背包系統 + 金幣經濟 + 出貨箱",
]
b = add_bullet_slide(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(5), items_left, font_size=16, color=LIGHT_GRAY, spacing=Pt(10))

# Right side: key features card
add_shape(slide, Inches(7.2), Inches(1.8), Inches(5.5), Inches(2.5), RGBColor(0x2D, 0x3F, 0x54))
add_text_box(slide, Inches(7.5), Inches(1.9), Inches(5), Inches(0.5), "✨ 核心功能", font_size=20, bold=True, color=ACCENT)
features = [
    "✔ 玩家角色繪製與面向指示",
    "✔ 兩種地形：草地（可通行／可耕作）、圍牆（障礙物）",
    "✔ 六階段農場狀態機 (FarmState)",
    "✔ 時間驅動的作物成長系統",
    "✔ 即時操作提示 (Toast)",
]
add_bullet_slide(slide, Inches(7.5), Inches(2.5), Inches(5), Inches(2), features, font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# Code stats
add_shape(slide, Inches(7.2), Inches(4.6), Inches(5.5), Inches(2.3), RGBColor(0x2D, 0x3F, 0x54))
add_text_box(slide, Inches(7.5), Inches(4.7), Inches(5), Inches(0.5), "📊 專案統計", font_size=20, bold=True, color=ACCENT3)
stats = [
    "• 7 個 Java 類別檔案",
    "• 總計約 690 行程式碼",
    "• 套件結構：farmgame",
    "• 純 JDK 開發，無外部依賴",
]
add_bullet_slide(slide, Inches(7.5), Inches(5.3), Inches(5), Inches(1.5), stats, font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

footer(slide)

# ================================================================
# SLIDE 4 — 開發環境與工具
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "開發環境與工具", "Development Environment")

env_items = [
    "語言：Java 21 (JDK 21 LTS, Eclipse Temurin)",
    "圖形 API：Swing / AWT（標準 JDK 函式庫）",
    "IDE：VS Code / 命令列",
    "作業系統：Windows 11",
    "建置方式：javac 直接編譯（無 Maven／Gradle）",
    "啟動方式：java -cp bin farmgame.Main",
]
add_bullet_slide(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(3), env_items, font_size=18, color=LIGHT_GRAY, spacing=Pt(10))

# 專案結構
add_shape(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(2.5), RGBColor(0x2D, 0x3F, 0x54))
add_text_box(slide, Inches(0.9), Inches(4.6), Inches(6), Inches(0.5), "📁 專案目錄結構", font_size=18, bold=True, color=ACCENT2)

tree = [
    "java/",
    " ├── src/farmgame/",
    " │   ├── Main.java          — 遊戲入口（JFrame）",
    " │   ├── GamePanel.java     — 遊戲迴圈 + 輸入 + UI",
    " │   ├── Player.java        — 玩家角色與碰撞",
    " │   ├── Tile.java          — 網格單元 + 狀態機",
    " │   ├── GameMap.java       — 二維地圖 + 成長邏輯",
    " │   ├── Inventory.java     — 背包（金幣／作物）",
    " │   └── ShippingBin.java   — 出貨箱",
    " ├── bin/                   — 編譯後 bytecode",
    " ├── compile.bat",
    " └── run.bat",
]
add_bullet_slide(slide, Inches(0.9), Inches(5.15), Inches(11), Inches(1.8), tree, font_size=13, color=LIGHT_GRAY, spacing=Pt(1))
footer(slide)

# ================================================================
# SLIDE 5 — 系統架構
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "系統架構", "System Architecture")

# Architecture diagram - boxes with arrows
add_shape(slide, Inches(3.8), Inches(1.8), Inches(1.8), Inches(1.0), ACCENT2)
add_text_box(slide, Inches(3.8), Inches(2.0), Inches(1.8), Inches(0.6), "Main.java\n入口", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.9), Inches(1.8), Inches(5.0), Inches(1.0), ACCENT)
add_text_box(slide, Inches(5.9), Inches(2.0), Inches(5.0), Inches(0.6), "GamePanel.java (遊戲主面板)", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.9), Inches(2.3), Inches(5.0), Inches(0.5), "遊戲迴圈 · 按鍵輸入 · 邏輯更新 · 渲染", font_size=10, color=RGBColor(0x1E, 0x2A, 0x38), alignment=PP_ALIGN.CENTER)

# Sub-modules
modules = [
    (Inches(0.5), Inches(3.5), Inches(2.6), Inches(1.3), "Tile.java", "地形列舉 + 農場狀態機", ACCENT2),
    (Inches(3.5), Inches(3.5), Inches(2.6), Inches(1.3), "Player.java", "網格移動 + 碰撞偵測 + 繪製", ACCENT2),
    (Inches(6.5), Inches(3.5), Inches(2.6), Inches(1.3), "GameMap.java", "二維地圖 + 成長更新 + 渲染", ACCENT2),
    (Inches(9.5), Inches(3.5), Inches(2.6), Inches(1.3), "Inventory.java", "背包 + 金幣", ACCENT3),
]
for left, top, w, h, title, desc, color in modules:
    add_shape(slide, left, top, w, h, color)
    add_text_box(slide, left, top + Inches(0.15), w, Inches(0.4), title, font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.55), w - Inches(0.2), Inches(0.6), desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ShippingBin inline
add_shape(slide, Inches(9.5), Inches(5.1), Inches(2.6), Inches(0.9), ACCENT3)
add_text_box(slide, Inches(9.5), Inches(5.2), Inches(2.6), Inches(0.4), "ShippingBin.java", font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.5), Inches(5.5), Inches(2.6), Inches(0.4), "出貨箱", font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Description text
add_text_box(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(1.2),
    "▸ 各類別均為 package farmgame 下的獨立檔案，職責單一\n▸ GamePanel 持有所有遊戲物件的參考，作為中央調度者\n▸ 渲染層次：地圖 → 出貨箱 → 玩家 → UI 疊層",
    font_size=14, color=LIGHT_GRAY)

footer(slide)

# ================================================================
# SLIDE 6 — Main.java
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "核心類別：Main.java", "遊戲入口 — JFrame 視窗啟動")

code = [
    "package farmgame;",
    "import javax.swing.*;",
    "",
    "public class Main {",
    "    public static void main(String[] args) {",
    "        try {",
    "            UIManager.setLookAndFeel(",
    "                UIManager.getSystemLookAndFeelClassName());",
    "        } catch (Exception ignored) {}",
    "",
    "        JFrame frame = new JFrame(\"農場物語 - 簡易農場遊戲\");",
    "        frame.setDefaultCloseOperation(JFrame.EXIT_ON_CLOSE);",
    "        frame.setResizable(false);",
    "        GamePanel panel = new GamePanel();",
    "        frame.add(panel);",
    "        frame.pack();",
    "        frame.setLocationRelativeTo(null);",
    "        frame.setVisible(true);",
    "    }",
    "}",
]
add_code_block(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(4.5), code, font_size=11)

add_text_box(slide, Inches(6.5), Inches(1.8), Inches(6), Inches(5),
    "職責：\n\n"
    "• 建立主視窗 JFrame，標題為「農場物語 - 簡易農場遊戲」\n\n"
    "• 設定系統原生外觀 (System Look & Feel)\n\n"
    "• 建立 GamePanel 實例並加入視窗\n\n"
    "• pack() 自動調整視窗大小為 800×600\n\n"
    "• setLocationRelativeTo(null) 將視窗置中\n\n"
    "• 視窗不可縮放 (setResizable(false))\n\n"
    "特點：\n"
    "• 僅 25 行，簡潔的入口點\n"
    "• 無外部依賴，純 JDK 即可編譯執行",
    font_size=14, color=LIGHT_GRAY)

footer(slide)

# ================================================================
# SLIDE 7 — GamePanel.java (Game Loop)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "核心類別：GamePanel.java", "遊戲主迴圈 (Game Loop) — 60 FPS 固定幀率")

code_loop = [
    "public class GamePanel extends JPanel implements Runnable {",
    "    private Thread gameThread;",
    "    private static final int TARGET_FPS = 60;",
    "    private static final long FRAME_MS = 1000 / TARGET_FPS;",
    "",
    "    private void startGameLoop() {",
    "        running = true;",
    "        gameThread = new Thread(this, \"GameLoop\");",
    "        gameThread.start();",
    "    }",
    "",
    "    @Override",
    "    public void run() {",
    "        long lastTime = System.nanoTime();",
    "        while (running) {",
    "            long now   = System.nanoTime();",
    "            long delta = (now - lastTime) / 1_000_000;",
    "            lastTime   = now;",
    "            if (delta > 200) delta = 200;",
    "            update(delta);  // 邏輯更新",
    "            repaint();      // 請求重繪",
    "            long elapsed = (System.nanoTime()-now)/1_000_000;",
    "            long sleep = FRAME_MS - elapsed;",
    "            if (sleep > 0) Thread.sleep(sleep);",
    "        }",
    "    }",
    "}",
]
add_code_block(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8), code_loop, font_size=11)

add_text_box(slide, Inches(6.8), Inches(1.8), Inches(6), Inches(5.2),
    "設計要點：\n\n"
    "❶ 獨立執行緒\n"
    "   遊戲邏輯不在 EDT 執行，避免阻塞 UI\n\n"
    "❷ Delta Time 機制\n"
    "   以 System.nanoTime() 計算每幀時間差\n"
    "   使成長計時不受 FPS 波動影響\n\n"
    "❸ 幀率控制\n"
    "   目標 60 FPS，每幀 ≈ 16.67 ms\n"
    "   計算本次耗時後補眠 (Thread.sleep)\n\n"
    "❹ 防跳幀\n"
    "   if (delta > 200) delta = 200\n"
    "   避免長時間卡頓後瞬間推進多個階段\n\n"
    "❺ 更新流程\n"
    "   update(delta) → 移動、工具、成長\n"
    "   repaint()     → paintComponent 繪製",
    font_size=14, color=LIGHT_GRAY)

footer(slide)

# ================================================================
# SLIDE 8 — GamePanel.java (Input Handling)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "核心類別：GamePanel.java", "按鍵輸入處理 (KeyListener + Edge Trigger)")

code_input = [
    "// 連續按壓（移動）",
    "private boolean upHeld, downHeld, leftHeld, rightHeld;",
    "// 單次動作（邊緣觸發）",
    "private boolean actionPressed, actionPrev;",
    "",
    "keyPressed(KeyEvent e) {",
    "    switch (e.getKeyCode()) {",
    "        case VK_W -> { upHeld = true; direction = 0; }",
    "        case VK_J -> actionPressed = true;",
    "        case VK_K -> seedPressed = true;",
    "        // ...",
    "    }",
    "}",
    "",
    "// Edge-triggered update block",
    "if (actionPressed && !actionPrev) useAction();",
    "if (seedPressed   && !seedPrev)   useSeeds();",
    "if (waterPressed  && !waterPrev)  useWateringCan();",
    "actionPrev = actionPressed;",
    "seedPrev   = seedPressed;",
    "// ...",
]
add_code_block(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5), code_input, font_size=11)

add_text_box(slide, Inches(6.8), Inches(1.8), Inches(6), Inches(5.2),
    "設計要點：\n\n"
    "❶ 雙模式按鍵\n"
    "   • 移動鍵 (WASD)：boolean held\n"
    "     持續按壓時為 true，釋放時為 false\n\n"
    "   • 動作鍵 (J/K/L/E)：edge-trigger\n"
    "     只在「按下瞬間」觸發一次\n\n"
    "❷ 邊緣觸發 (Edge Trigger)\n"
    "   使用 _Pressed / _Prev 配對：\n"
    "   if (actionPressed && !actionPrev) → 只執行一次\n"
    "   避免 OS 按鍵重複觸發多次動作\n\n"
    "❸ 移動冷卻 (Move Cooldown)\n"
    "   MOVE_INTERVAL_MS = 130ms\n"
    "   每 130ms 才可移動一格\n"
    "   防止按住時 60fps 瞬間飛越整個地圖\n\n"
    "❹ Toast 提示系統\n"
    "   每次操作顯示 2.5 秒訊息回饋",
    font_size=14, color=LIGHT_GRAY)

footer(slide)

# ================================================================
# SLIDE 9 — Player.java
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "核心類別：Player.java", "玩家角色 — 網格移動與視覺繪製")

code_player = [
    "public class Player {",
    "    public int gridRow, gridCol;  // 網格座標",
    "    public int direction = 1;     // 0=上 1=下 2=左 3=右",
    "",
    "    public void move(int dRow, int dCol) {",
    "        int tr = gridRow + dRow;",
    "        int tc = gridCol + dCol;",
    "        if (gameMap.isWalkable(tr, tc)) {",
    "            gridRow = tr;",
    "            gridCol = tc;",
    "        }",
    "    }",
    "",
    "    public void render(Graphics g) {",
    "        // 藍色襯衫 + 膚色頭部 + 黑色眼睛",
    "        // 面向箭頭（白色三角形）",
    "        // 除錯：黃色邊框高亮所在格",
    "    }",
    "}",
]
add_code_block(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(3.5), code_player, font_size=12)

add_text_box(slide, Inches(6.5), Inches(1.8), Inches(6), Inches(5.2),
    "設計要點：\n\n"
    "❶ 網格座標系統\n"
    "   gridRow / gridCol 定位在 25×19 網格中\n"
    "   非像素座標，簡化碰撞與互動計算\n\n"
    "❷ Grid-based 移動\n"
    "   一次移動一整格 (32×32 px)\n"
    "   搭配 GameMap.isWalkable() 檢查目標格\n\n"
    "❸ 碰撞偵測\n"
    "   isWalkable = TileType != WALL\n"
    "   越界自動回傳 false\n\n"
    "❹ 面向系統\n"
    "   direction 影響工具使用的目標格\n"
    "   渲染時以白色小箭頭指示面向\n\n"
    "❺ 繪製層次\n"
    "   頭部 → 身體 → 眼睛 → 面向箭頭\n"
    "   全部使用純 AWT Graphics2D 繪圖",
    font_size=14, color=LIGHT_GRAY)

footer(slide)

# ================================================================
# SLIDE 10 — Tile.java (State Machine)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "核心類別：Tile.java", "網格單元 — 地形列舉 + 農場狀態機 (State Machine)")

code_tile = [
    "public class Tile {",
    "    public enum TileType {",
    "        GRASS,  // 草地（可通行、可耕作）",
    "        WALL    // 障礙物（不可通行）",
    "    }",
    "",
    "    public enum FarmState {",
    "        NONE,    // 一般草地",
    "        TILLED,  // 已翻土",
    "        PLANTED, // 已播種",
    "        WATERED, // 已澆水",
    "        GROWING, // 成長中（4階段）",
    "        MATURE   // 成熟可收成",
    "    }",
    "",
    "    public TileType type;",
    "    public FarmState farmState = NONE;",
    "    public String cropType = \"\";",
    "    public int growthStage = 0;",
    "    public int growthTimerMs = 0;",
    "",
    "    static final int TIME_WATERED_TO_GROWING = 10000;",
    "    static final int TIME_PER_GROWTH_STAGE = 2500;",
    "    static final int TOTAL_GROWTH_STAGES = 4;",
    "}",
]
add_code_block(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5), code_tile, font_size=11)

add_text_box(slide, Inches(6.8), Inches(1.8), Inches(6), Inches(5.2),
    "設計要點：\n\n"
    "❶ 雙列舉結構\n"
    "   TileType：定義地形本質（不變）\n"
    "   FarmState：定義農場狀態（可變）\n"
    "   分離「是什麼」與「變成什麼」\n\n"
    "❷ 六階段狀態機\n"
    "   NONE → TILLED → PLANTED → WATERED\n"
    "   → GROWING(stage 1~4) → MATURE\n"
    "   每個階段僅接受特定輸入轉換\n\n"
    "❸ 時間驅動成長\n"
    "   • 澆水後 10 秒萌芽 (GROWING stage 1)\n"
    "   • 每 2.5 秒成長一個階段\n"
    "   • 4 階段後成熟 (MATURE)\n\n"
    "❹ 擴充性\n"
    "   可輕易新增作物種類、\n"
    "   不同成長時間、更多地形類型",
    font_size=14, color=LIGHT_GRAY)

footer(slide)

# ================================================================
# SLIDE 11 — GameMap.java
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "核心類別：GameMap.java", "二維地圖 — 生成、成長、渲染")

code_map = [
    "public class GameMap {",
    "    static final int TILE_SIZE = 32;",
    "    static final int COLS = 25, ROWS = 19;",
    "    private final Tile[][] tiles;",
    "",
    "    public GameMap() {",
    "        // 邊界 WALL，內部 GRASS",
    "        tiles = new Tile[ROWS][COLS];",
    "        for (int r...) for (int c...)",
    "            tiles[r][c] = new Tile(",
    "                (r==0||r==ROWS-1||c==0||c==COLS-1)",
    "                ? WALL : GRASS);",
    "    }",
    "",
    "    public void updateGrowth(long deltaMs) {",
    "        for each tile:",
    "          if WATERED+crop → timer+=delta,",
    "             if timer≥10s → GROWING stage1",
    "          if GROWING → timer+=delta,",
    "             if timer≥2.5s → next stage",
    "    }",
    "",
    "    public void render(Graphics g) {",
    "        for each tile:",
    "          draw base (WALL/GRASS)",
    "          if farmState != NONE: draw overlay",
    "          draw grid lines",
    "    }",
    "}",
]
add_code_block(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5), code_map, font_size=11)

add_text_box(slide, Inches(6.8), Inches(1.8), Inches(6), Inches(5.2),
    "設計要點：\n\n"
    "❶ 二維陣列 (Tile[19][25])\n"
    "   共 475 格，每格 32×32 像素\n\n"
    "❷ 地圖生成\n"
    "   邊界自動設為 WALL（圍牆）\n"
    "   內部全部 GRASS（草地）\n\n"
    "❸ 成長系統 (updateGrowth)\n"
    "   遍歷每格，根據 FarmState 累計時間\n"
    "   時間到則轉換狀態 / 推進階段\n\n"
    "❹ 渲染層次\n"
    "   基底（WALL 棕色 / GRASS 綠色）\n"
    "   → 農場疊層（犁溝、種子、水光、作物）\n"
    "   → 極淡網格線\n\n"
    "❺ 視覺細節\n"
    "   • 草地顏色略有隨機變化\n"
    "   • 翻土有三條犁溝紋理\n"
    "   • 澆水有水光高亮\n"
    "   • 作物 5 階段視覺（含果實）",
    font_size=14, color=LIGHT_GRAY)

footer(slide)

# ================================================================
# SLIDE 12 — Inventory + ShippingBin
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "核心類別：Inventory & ShippingBin", "背包系統與經濟循環")

code_inv = [
    "// Inventory.java — 背包",
    "public class Inventory {",
    "    public int gold = 0;",
    "    public int cropCount = 0;",
    "",
    "    public void harvestCrop() { cropCount++; }",
    "",
    "    public int sellAll() {",
    "        int sold = cropCount;",
    "        if (sold > 0) {",
    "            gold += sold * 10;  // 每個 10 G",
    "            cropCount = 0;",
    "        }",
    "        return sold;",
    "    }",
    "}",
]
add_code_block(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(3.2), code_inv, font_size=11)

code_bin = [
    "// ShippingBin.java — 出貨箱",
    "public class ShippingBin {",
    "    public final int row, col;",
    "    // 繪製木箱本體 + 箱蓋 + 木紋 + 標籤",
    "}",
]
add_code_block(slide, Inches(0.6), Inches(5.2), Inches(5.5), Inches(1.5), code_bin, font_size=12)

add_text_box(slide, Inches(6.5), Inches(1.8), Inches(6), Inches(5.2),
    "Inventory 設計：\n\n"
    "• 簡單的計數器：gold 與 cropCount\n"
    "• harvestCrop() 收成時 ++ 計數\n"
    "• sellAll() 將所有作物賣出，每個 10 G\n\n"
    "ShippingBin 設計：\n\n"
    "• 固定於地圖 (COLS-3, ROWS-3)\n"
    "• 使用 AWT Graphics 繪製木箱視覺\n"
    "• 玩家在相鄰格按 E 觸發交易\n\n"
    "經濟循環：\n"
    "  耕作 → 收成 → 出貨箱賣出\n"
    "  → 獲得金幣 → （未來）購買種子/工具\n\n"
    "互動流程 (GamePanel.interact)：\n"
    "  if 曼哈頓距離 ≤ 1:\n"
    "     inventory.sellAll()\n"
    "     showToast(\"賣出 X 個作物...\")",
    font_size=14, color=LIGHT_GRAY)

footer(slide)

# ================================================================
# SLIDE 13 — 耕作流程圖
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "耕作流程與狀態機", "Farming State Machine — 完整生命週期")

# Flow diagram using shapes
states = [
    (Inches(0.4), Inches(2.2), Inches(1.8), Inches(1.0), "草地\nGRASS", RGBColor(0x4E, 0xCB, 0x71)),
    (Inches(2.6), Inches(2.2), Inches(1.8), Inches(1.0), "翻土\nTILLED", RGBColor(0x8B, 0x5A, 0x2B)),
    (Inches(4.8), Inches(2.2), Inches(1.8), Inches(1.0), "播種\nPLANTED", RGBColor(0x78, 0x50, 0x28)),
    (Inches(7.0), Inches(2.2), Inches(1.8), Inches(1.0), "澆水\nWATERED", RGBColor(0x46, 0x41, 0x55)),
    (Inches(9.2), Inches(2.2), Inches(1.8), Inches(1.0), "成長中\nGROWING", RGBColor(0x32, 0xB4, 0x32)),
    (Inches(11.4), Inches(2.2), Inches(1.6), Inches(1.0), "成熟\nMATURE", RGBColor(0xFF, 0xD7, 0x00)),
]
for left, top, w, h, label, color in states:
    add_shape(slide, left, top, w, h, color)
    add_text_box(slide, left, top + Inches(0.2), w, Inches(0.5), label, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Arrow labels
arrows = [
    (Inches(2.1), Inches(2.0), "J 鋤頭", ACCENT),
    (Inches(4.3), Inches(2.0), "K 播種", ACCENT2),
    (Inches(6.5), Inches(2.0), "L 澆水", ACCENT2),
    (Inches(8.7), Inches(2.0), "10 秒", ACCENT3),
    (Inches(10.9), Inches(2.0), "每 2.5s\n×4 階段", ACCENT3),
]
for left, top, label, color in arrows:
    add_text_box(slide, left, top, Inches(1.6), Inches(0.8), label, font_size=9, bold=True, color=color, alignment=PP_ALIGN.CENTER)

# Harvest arrow back
add_text_box(slide, Inches(2.6), Inches(3.5), Inches(1.8), Inches(0.9), "J 收成\n(循環)", font_size=11, bold=True, color=RGBColor(0xFF, 0x6B, 0x6B), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.0), Inches(3.5), Inches(0.3), Inches(0.3), "↻", font_size=20, bold=True, color=RGBColor(0xFF, 0x6B, 0x6B))

# Detailed description
details = [
    "狀態轉換說明：",
    "",
    "▸ 僅 GRASS + NONE 狀態可接受「J 鋤頭」→ 變成 TILLED",
    "▸ 僅 TILLED 狀態可接受「K 種子」→ 變成 PLANTED（作物：防風草）",
    "▸ TILLED 或 PLANTED 可接受「L 澆水壺」→ 變成 WATERED",
    "▸ WATERED + 有作物 → 等待 10 秒 → GROWING stage 1",
    "▸ GROWING → 每 2.5 秒推進一階段 → stage 4 後 → MATURE",
    "▸ MATURE 可接受「J 收成」→ 作物進背包，變回 TILLED",
]
add_bullet_slide(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(3.5), details, font_size=14, color=LIGHT_GRAY, spacing=Pt(3))

footer(slide)

# ================================================================
# SLIDE 14 — 操作說明
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "操作說明", "Controls Reference")

keys = [
    ("W A S D", "移動角色（上／下／左／右）", ACCENT2),
    ("J", "鋤頭（草地 → 翻土）／收成（成熟作物 → 背包）", ACCENT),
    ("K", "種子（翻土 → 播種）", ACCENT),
    ("L", "澆水壺（播種後 → 澆水）", ACCENT),
    ("E", "互動（靠近出貨箱 → 賣出所有作物）", ACCENT3),
]

for i, (key, desc, color) in enumerate(keys):
    y = Inches(2.0) + i * Inches(0.85)
    # Key badge
    shape = add_shape(slide, Inches(1.0), y, Inches(2.0), Inches(0.6), color)
    shape.text_frame.paragraphs[0].text = key
    shape.text_frame.paragraphs[0].font.size = Pt(16)
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.name = 'Consolas'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].space_before = Pt(4)
    # Description
    add_text_box(slide, Inches(3.3), y + Inches(0.05), Inches(9), Inches(0.5), desc, font_size=16, color=LIGHT_GRAY)

# Hint
add_text_box(slide, Inches(1.0), Inches(6.2), Inches(11), Inches(0.8),
    "💡 提示：操作成功時畫面中央會顯示 Toast 提示訊息（2.5 秒）\n   靠近出貨箱（右下角木箱）按 E 可將背包作物變現",
    font_size=14, color=SOFT_GREEN)

footer(slide)

# ================================================================
# SLIDE 15 — 成果展示
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "成果展示", "Screenshots & Demo")

# Left card - visual features
add_shape(slide, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.2), RGBColor(0x2D, 0x3F, 0x54))
add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.5), "🖼 畫面渲染層次", font_size=20, bold=True, color=ACCENT)

layers = [
    "1. 地圖底層（WALL 棕色石牆 + GRASS 綠色草地）",
    "2. 農場狀態疊層（翻土棕色、播種深褐、澆水深紫）",
    "3. 作物繪製（4 階段莖葉 + 成熟金黃果實）",
    "4. 出貨箱（木箱視覺 + 文字標籤）",
    "5. 玩家角色（藍襯衫 + 膚色頭 + 面向箭頭）",
    "6. UI 疊層（半透明頂欄 + 金幣/背包/操作提示）",
    "7. Toast 訊息（畫面中央 2.5 秒）",
]
add_bullet_slide(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4), layers, font_size=14, color=LIGHT_GRAY, spacing=Pt(6))

# Right card - tech highlights
add_shape(slide, Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.2), RGBColor(0x2D, 0x3F, 0x54))
add_text_box(slide, Inches(7.1), Inches(1.9), Inches(5.5), Inches(0.5), "⚡ 技術亮點", font_size=20, bold=True, color=ACCENT3)

techs = [
    "✅ 固定 60 FPS 遊戲迴圈 + Delta Time 機制",
    "✅ 邊緣觸發 (Edge Trigger) 按鍵系統",
    "✅ Grid-based 碰撞偵測（不可穿越圍牆）",
    "✅ 狀態機驅動的耕作系統（6 階段）",
    "✅ 時間驅動作物成長（10s→萌芽，2.5s/階段）",
    "✅ 背包 + 金幣經濟循環",
    "✅ 純 AWT Graphics 繪製，無外部資源",
    "✅ 模組化類別設計，職責單一",
    "✅ 中文註解 + Toast 即時回饋",
]
add_bullet_slide(slide, Inches(7.1), Inches(2.5), Inches(5.5), Inches(4), techs, font_size=14, color=LIGHT_GRAY, spacing=Pt(6))

footer(slide)

# ================================================================
# SLIDE 16 — 未來展望
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "未來展望", "Future Roadmap")

features = [
    "🌱 更多作物種類（番茄、草莓、南瓜等，各自不同成長週期與售價）",
    "🏗 更多工具（斧頭砍樹、十字鎬挖礦、釣竿釣魚）",
    "🗺 更大多元地圖（森林、湖泊、礦區、村民房屋）",
    "👨‍🌾 NPC 村民系統（對話、送禮、好感度）",
    "📅 季節與天氣系統（春夏秋冬 + 晴雨雪）",
    "🔊 音效與背景音樂",
    "🎮 滑鼠支援（點擊移動、點擊耕作）",
    "💾 存檔／讀檔（Serialization 或 JSON）",
    "🖼 圖片資源取代程式繪圖（Sprite / TileSet）",
    "🌐 在地化支援（正體中文／英文切換）",
]
add_bullet_slide(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(5), features, font_size=16, color=LIGHT_GRAY, spacing=Pt(8))

footer(slide)

# ================================================================
# SLIDE 17 — Thank You
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(2.5), prs.slide_width, Inches(2.5), BG_CARD)
add_shape(slide, Inches(0), Inches(2.2), prs.slide_width, Inches(0.08), ACCENT)

add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.0), "感謝聆聽", font_size=48, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6), "Thank You", font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5), "以 Java Swing/AWT 實作之 2D 簡易農場遊戲", font_size=14, color=SOFT_GREEN, alignment=PP_ALIGN.CENTER)

footer(slide, "")

# ================================================================
# SAVE
# ================================================================
output_path = "C:\\Users\\user\\Desktop\\java\\FarmGame_Report.pptx"
prs.save(output_path)
print(f"PPT saved to {output_path}")
